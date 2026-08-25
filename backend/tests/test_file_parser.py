"""Extractor coverage for file_parser (M7 upload pipeline).

Why this file exists: the naive extractors only read python-docx `doc.paragraphs`
and top-level pptx text frames, so table-heavy docx and grouped/table pptx —
the norm for real course materials — extracted as "" and never reached the
knowledge base. These tests pin the fixed behavior: paragraphs + tables in
document order for DOCX; recursive group shapes + tables + notes for PPTX.
"""
import io
import sys
import unittest
from pathlib import Path

_BACKEND = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(_BACKEND))

from app.core.file_parser import extract_text  # noqa: E402


def _docx_bytes(paragraphs=(), tables=()):
    """Build a .docx in memory: interleaved paragraphs and 2x2 tables."""
    import docx
    doc = docx.Document()
    for item in paragraphs:
        doc.add_paragraph(item)
    for rows in tables:
        t = doc.add_table(rows=len(rows), cols=len(rows[0]))
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                t.cell(i, j).text = val
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _pptx_bytes(build):
    """Build a .pptx in memory; `build(slide, shapes)` adds content."""
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    build(slide)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestDocxExtraction(unittest.TestCase):
    def test_paragraphs_still_work(self):
        raw = _docx_bytes(paragraphs=["光的折射定义", "折射定律 n1sinθ1=n2sinθ2"])
        text = extract_text("a.docx", raw)
        self.assertIn("光的折射定义", text)
        self.assertIn("n1sinθ1", text)

    def test_table_content_extracted(self):
        raw = _docx_bytes(tables=[[["知识点", "浮力"], ["定义", "向上的托力"]]])
        text = extract_text("t.docx", raw)
        self.assertIn("浮力", text)
        self.assertIn("向上的托力", text)
        self.assertIn("|", text)  # cells joined as "cell | cell"

    def test_table_only_docx_not_empty(self):
        # Regression: table-only documents used to extract as "".
        raw = _docx_bytes(tables=[[["串并联", "电流规律"]]])
        self.assertTrue(extract_text("only.docx", raw).strip())

    def test_mixed_order_paragraph_then_table(self):
        raw = _docx_bytes(paragraphs=["前置段落"], tables=[[["表内文字", "x"]]])
        text = extract_text("m.docx", raw)
        self.assertLess(text.index("前置段落"), text.index("表内文字"))


class TestPptxExtraction(unittest.TestCase):
    def test_top_level_textbox_still_works(self):
        from pptx.util import Inches

        def build(slide):
            tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
            tb.text_frame.text = "勾股定理 a^2+b^2=c^2"
        text = extract_text("a.pptx", _pptx_bytes(build))
        self.assertIn("勾股定理", text)

    def test_group_shape_text_extracted(self):
        from pptx.util import Inches

        def build(slide):
            grp = slide.shapes.add_group_shape()
            tb = grp.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
            tb.text_frame.text = "组合形状里的惯性定律"
        text = extract_text("g.pptx", _pptx_bytes(build))
        self.assertIn("惯性定律", text)

    def test_table_cells_extracted(self):
        from pptx.util import Inches

        def build(slide):
            tbl = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(1)).table
            tbl.cell(0, 0).text = "电路"
            tbl.cell(1, 1).text = "处处相等"
        text = extract_text("t.pptx", _pptx_bytes(build))
        self.assertIn("电路", text)
        self.assertIn("处处相等", text)

    def test_notes_extracted(self):
        from pptx.util import Inches

        def build(slide):
            tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
            tb.text_frame.text = "正文"
            slide.notes_slide.notes_text_frame.text = "备注：阿基米德原理补充"
        text = extract_text("n.pptx", _pptx_bytes(build))
        self.assertIn("阿基米德原理", text)

    def test_slide_boundary_formfeed_kept(self):
        from pptx.util import Inches
        from pptx import Presentation
        prs = Presentation()
        for label in ("第一页", "第二页"):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
            tb.text_frame.text = label
        buf = io.BytesIO()
        prs.save(buf)
        text = extract_text("s.pptx", buf.getvalue())
        self.assertIn("\f", text)  # page boundary feeds the chunker's page metadata


class TestTextExtraction(unittest.TestCase):
    def test_txt_utf8_and_gb18030(self):
        self.assertIn("浮力", extract_text("a.txt", "浮力定义".encode("utf-8")))
        self.assertIn("浮力", extract_text("a.txt", "浮力定义".encode("gb18030")))

    def test_md_read_as_text(self):
        self.assertIn("# 标题", extract_text("n.md", "# 标题".encode("utf-8")))


if __name__ == "__main__":
    unittest.main()
