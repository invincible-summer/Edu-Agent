"""Extract text from uploaded course materials: PDF / DOCX / PPTX / TXT / MD.

Each parser is isolated so a failure in one format never breaks the others.
Returns "" on any extraction failure (the caller treats that as no-text).
"""
from __future__ import annotations

from pathlib import Path


def extract_text(filename: str, raw: bytes, *, ocr_fallback: bool = True) -> str:
    """Dispatch by extension. Returns extracted plain text or "".

    ``ocr_fallback``: 扫描 PDF 文本层稀疏时同步 tesseract OCR 回退（对话/资料库
    上传用，受 PDF_OCR_SYNC_MAX_PAGES 保护）。教材库上传传 False——它的扫描 PDF
    由 textbook_builder 后台 async OCR（视觉模型优先），避免同步段重复 OCR。
    """
    lower = (filename or "").lower()
    try:
        if lower.endswith(".pdf"):
            return _extract_pdf(raw, ocr_fallback=ocr_fallback)
        if lower.endswith(".docx"):
            return _extract_docx(raw)
        if lower.endswith(".pptx"):
            return _extract_pptx(raw)
        if lower.endswith((".txt", ".md", ".markdown")):
            return _extract_text_bytes(raw)
    except Exception:
        return ""
    return ""


def _extract_pdf(raw: bytes, *, ocr_fallback: bool = True) -> str:
    import fitz
    from .pdf_ocr import FITZ_LOCK  # PyMuPDF 非线程安全：文档操作全程持锁
    with FITZ_LOCK:
        doc = fitz.open(stream=raw, filetype="pdf")
        # \f (form feed) joins pages: the structure-aware chunker uses it as a
        # hard page boundary so chunks keep their page number for citation.
        pages = [page.get_text() for page in doc]
        n = doc.page_count
        doc.close()
    text = "\f".join(pages)
    if not ocr_fallback:
        return text
    # 扫描/混合版回退：存在稀疏页（扫描页/图片页）→ 同步 tesseract 逐页择优 OCR
    # （仅本地、受 PDF_OCR_SYNC_MAX_PAGES 保护，避免大书卡住上传请求；视觉模型通道
    # 只在教材库后台 async 用）。稀疏页用 OCR 文本、达标页保留文本层，按页合并。
    ocr_text = _maybe_sync_ocr(raw, text, n, pages=pages)
    if ocr_text:
        return ocr_text
    return text


def _maybe_sync_ocr(raw: bytes, text_layer: str, n_pages: int,
                    *, pages: list[str] | None = None) -> str:
    """同步 tesseract OCR 回退（扫描/混合 PDF）。返回合并文本或 ""（不触发/失败）。

    触发条件（P5a-A2 逐页判定 + P8 乱码页，替代旧的全书页均判定）：
      - mode=off：不触发。
      - mode=auto：存在**任一**稀疏页（页字符 < 阈值）或稠密乱码页（定制字体
        无 ToUnicode 映射，见 core/text_quality）即触发，但只 OCR 这些页，
        文本层达标页原样保留（混合书两半都得最佳文本；良好文本层页不被降质）。
      - mode=on：强制整本 OCR。
    tesseract 缺失时放弃（教材库后台仍有视觉模型通道）。OCR 页数受
    PDF_OCR_SYNC_MAX_PAGES 保护；空页以 "" 占位拼接（A1：页码与物理页对齐）。
    """
    try:
        from .config import settings
        from . import pdf_ocr
    except Exception:
        return ""
    if settings.pdf_ocr_mode == "off" or n_pages <= 0:
        return ""
    import shutil
    if not shutil.which("tesseract"):
        return ""  # tesseract 未安装：同步段无视觉模型通道，放弃（教材库后台仍可走视觉模型）
    from . import ocr as _ocr
    cap = settings.pdf_ocr_sync_max_pages
    if settings.pdf_ocr_mode == "on":
        # 强制整本 OCR（文本层乱码等场景）；cap 限制总页数。
        full = pdf_ocr.ocr_pdf_pages_sync(
            raw, lambda png: _ocr._tesseract_ocr(png, psm=3),
            max_pages=cap, dpi=settings.pdf_ocr_dpi)
        ocr_text = "\f".join(full)  # 空页占位保留（页码对齐）
        if not ocr_text.strip():
            return ""
        if n_pages > cap:
            ocr_text += f"\n[注：页数 {n_pages} 超过同步 OCR 上限 {cap}，仅识别前 {cap} 页；完整识别请上传到教材库]"
        return ocr_text
    # auto：逐页判定（稀疏 ∪ 乱码），无目标页不触发。
    page_texts = pages if pages is not None else text_layer.split("\f")
    if not pdf_ocr.pages_needing_ocr(page_texts):
        return ""
    merged, stats = pdf_ocr.ocr_pdf_pages_mixed_sync(
        raw, lambda png: _ocr._tesseract_ocr(png, psm=3),
        page_texts=page_texts, max_ocr_pages=cap, dpi=settings.pdf_ocr_dpi)
    if not merged:
        return ""
    ocr_text = "\f".join(merged)  # 空页占位保留（页码与物理页一一对应）
    if not ocr_text.strip():
        return ""
    if stats["sparse"] > stats["ocr_done"]:
        ocr_text += (f"\n[注：扫描/图片页共 {stats['sparse']} 页，超过同步 OCR 上限，"
                     f"仅识别前 {stats['ocr_done']} 页；完整识别请上传到教材库]")
    return ocr_text


def _extract_docx(raw: bytes) -> str:
    """Paragraphs AND tables, in document order.

    python-docx's `doc.paragraphs` skips table content entirely — and course
    materials are table-heavy — so we walk the body XML directly and expand
    each table row as "cell | cell" text (nested tables recursed).
    """
    import io
    import docx
    from docx.table import Table, _Cell
    from docx.text.paragraph import Paragraph

    doc = docx.Document(io.BytesIO(raw))

    def cell_text(cell: _Cell) -> str:
        parts = [p.text for p in cell.paragraphs if p.text.strip()]
        for tbl in cell.tables:  # nested tables
            parts.extend(table_lines(tbl))
        return " ".join(parts)

    def table_lines(tbl: Table) -> list[str]:
        lines = []
        for row in tbl.rows:
            line = " | ".join(t for t in (cell_text(c).strip() for c in row.cells) if t)
            if line:
                lines.append(line)
        return lines

    out: list[str] = []
    for child in doc.element.body.iterchildren():
        if child.tag.endswith("}p"):
            p = Paragraph(child, doc)
            if p.text.strip():
                out.append(p.text)
        elif child.tag.endswith("}tbl"):
            out.extend(table_lines(Table(child, doc)))
    return "\n".join(out)


def _extract_pptx(raw: bytes) -> str:
    """Slide text: text frames + tables + group shapes (recursive) + notes.

    The naive top-level `shape.has_text_frame` walk misses grouped shapes,
    table cells and speaker notes — all common in real course decks.
    """
    import io
    import pptx
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = pptx.Presentation(io.BytesIO(raw))

    def shape_texts(shape) -> list[str]:
        texts: list[str] = []
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                texts.extend(shape_texts(sub))
            return texts
        if shape.has_text_frame:
            texts.extend(run.text for run in shape.text_frame.paragraphs
                         for run in run.runs if run.text)
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                line = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                if line:
                    texts.append(line)
        return texts

    parts: list[str] = []
    for slide in prs.slides:
        texts: list[str] = []
        for shape in slide.shapes:
            texts.extend(shape_texts(shape))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(notes)
        if texts:
            parts.append("\n".join(texts))
    # \f between slides: same hard-boundary convention as PDF pages, so the
    # chunker records the slide number as the chunk's page metadata.
    return "\f".join(parts)


def _extract_text_bytes(raw: bytes) -> str:
    for enc in ("utf-8", "gb18030", "latin-1"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    return ""


SUPPORTED_EXTS = (".pdf", ".docx", ".pptx", ".txt", ".md", ".markdown")
MAX_UPLOAD_BYTES = 256 * 1024 * 1024  # 256 MB — 整本扫描教材可达数百 MB（P6-A1）
MAX_IMAGE_BYTES = 20 * 1024 * 1024  # 20 MB — inline chat images stay small
